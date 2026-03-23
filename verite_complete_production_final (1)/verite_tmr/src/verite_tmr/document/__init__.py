"""
Document ingestion — multi-format document processing and concept extraction.
Restored from v9 EnhancedDocProcessor + TMRAIEngine.
"""
from __future__ import annotations
import json, re, warnings
from pathlib import Path
from typing import Any
import numpy as np

class DocumentProcessor:
    """Multi-format document ingestion: PDF, DOCX, PPTX, TXT, MP3/MP4 (Whisper)."""
    def __init__(self, device="cpu"):
        self._device = device; self._whisper = None

    def extract_text(self, path: str) -> str:
        p = Path(path); ext = p.suffix.lower()
        dispatch = {".pdf": self._pdf, ".docx": self._docx, ".pptx": self._pptx,
                    ".txt": self._txt, ".md": self._txt,
                    ".mp3": self._audio, ".wav": self._audio, ".mp4": self._video}
        fn = dispatch.get(ext)
        if fn is None:
            raise ValueError(f"Unsupported format: {ext}")
        return fn(str(p))

    def _pdf(self, p):
        try:
            import pdfplumber
            text = []
            with pdfplumber.open(p) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t: text.append(t)
            if text: return "\n".join(text)
        except Exception: pass
        # OCR fallback
        try:
            from pdf2image import convert_from_path
            import pytesseract
            imgs = convert_from_path(p, dpi=200)
            return "\n".join(pytesseract.image_to_string(img) for img in imgs)
        except Exception as e:
            raise RuntimeError(f"PDF extraction failed: {e}")

    def _docx(self, p):
        from docx import Document
        doc = Document(p)
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

    def _pptx(self, p):
        from pptx import Presentation
        prs = Presentation(p); texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)

    def _txt(self, p):
        return Path(p).read_text(encoding="utf-8", errors="replace")

    def _audio(self, p):
        self._load_whisper()
        if self._whisper:
            result = self._whisper.transcribe(p)
            return result.get("text", "")
        return ""

    def _video(self, p):
        import subprocess, tempfile
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            subprocess.run(["ffmpeg","-i",p,"-ac","1","-ar","16000",tmp.name,"-y"],
                          capture_output=True, check=True)
            return self._audio(tmp.name)

    def _load_whisper(self):
        if self._whisper is None:
            try:
                import whisper
                self._whisper = whisper.load_model("base", device=self._device)
            except ImportError:
                warnings.warn("whisper not installed — audio transcription unavailable")

    def cleanup(self):
        self._whisper = None


class ConceptExtractor:
    """AI-powered concept extraction with Groq/Gemini/TF-IDF cascade."""
    def __init__(self, groq_key="", gemini_key=""):
        self._groq = self._gemini = None
        if groq_key:
            try:
                from groq import Groq
                self._groq = Groq(api_key=groq_key)
            except Exception: pass
        if gemini_key:
            try:
                from google import genai
                self._gemini = genai.Client(api_key=gemini_key)
            except Exception: pass

    def extract_concepts(self, text: str, max_concepts: int = 10) -> list[dict]:
        """Extract key concepts using LLM cascade or TF-IDF fallback."""
        # Try LLM first
        prompt = (f'Extract {max_concepts} key concepts for study.\n'
                  f'Return JSON array: [{{"concept":"name","definition":"explanation",'
                  f'"difficulty":"easy|medium|hard","category":"topic"}}]\n\nTEXT:\n{text[:200000]}')

        raw = self._llm_call(prompt)
        if raw:
            try:
                concepts = json.loads(raw.strip().strip("```json").strip("```"))
                if isinstance(concepts, list) and concepts:
                    return [c for c in concepts if "concept" in c and "definition" in c][:max_concepts]
            except Exception: pass
        return self._tfidf_concepts(text, max_concepts)

    def _llm_call(self, prompt):
        if self._groq:
            try:
                r = self._groq.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[{"role":"user","content":prompt}], max_tokens=4096)
                return r.choices[0].message.content
            except Exception: pass
        if self._gemini:
            try:
                r = self._gemini.models.generate_content(model="gemini-2.0-flash", contents=prompt)
                return r.text
            except Exception: pass
        return None

    def _tfidf_concepts(self, text, n):
        try:
            from sklearn.feature_extraction.text import TfidfVectorizer
            sents = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if len(s.strip())>30]
            if not sents: return []
            vec = TfidfVectorizer(max_features=500, stop_words='english', ngram_range=(1,3))
            mat = vec.fit_transform(sents)
            names = vec.get_feature_names_out()
            means = mat.mean(axis=0).A1
            top_idx = means.argsort()[::-1]
            concepts = []
            for idx in top_idx[:n]:
                term = names[idx]
                best_sent = sents[mat.getcol(idx).toarray().argmax()]
                concepts.append({"concept": term, "definition": best_sent[:200],
                                "difficulty": "medium", "category": "general"})
            return concepts
        except ImportError:
            return []

    def compute_similarity(self, concepts):
        try:
            from sklearn.feature_extraction.text import TfidfVectorizer
            from sklearn.metrics.pairwise import cosine_similarity
            texts = [f"{c['concept']} {c['definition']}" for c in concepts]
            if len(texts) < 2: return np.zeros((1,1))
            mat = TfidfVectorizer(stop_words='english').fit_transform(texts)
            sim = cosine_similarity(mat); np.fill_diagonal(sim, 0)
            return sim
        except: return np.zeros((len(concepts), len(concepts)))

__all__ = ["DocumentProcessor", "ConceptExtractor"]
